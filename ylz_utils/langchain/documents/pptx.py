from __future__ import annotations
from typing import TYPE_CHECKING
if TYPE_CHECKING:
    from ylz_utils.langchain import LangchainLib

#from langchain_community.document_loaders import UnstructuredPowerPointLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
import pptx 
from pptx.slide import Slide,SlideLayout,SlideShapes
from pptx.presentation import Presentation
from typing import List,Dict, Literal, Tuple
from pptx.shapes.autoshape import Shape, MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData,XyChartData,BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches,Cm,Pt
class PptxLib():
    ppt: Presentation  = None
    slides: List[Slide] = []
    unit = None
    def __init__(self,langchainLib:LangchainLib):
        self.langchainLib = langchainLib
        self.unit = Pt
    def set_default_unit(self,unit):
        if unit.lower()=="inches":
            self.unit = Inches
        elif unit.lower()=="cm":
            self.unit = Cm
        elif unit.lower()=="pt":
            self.unit = Pt
        else:
            raise Exception("必须指定Inches,Cm,Pt中的一种单位")
    def _get_sizes(self,left:float|int,top:float|int,width:float|int,height:float|int):
        if self.unit == Inches:
            return Inches(left),Inches(top),Inches(width),Inches(height)
        elif self.unit == Cm:
            return Cm(left),Cm(top),Cm(width),Cm(height)
        else:
            return Pt(left),Pt(top),Pt(width),Pt(height)
    def _get_size(self,size:float|int):
        if self.unit == Inches:
            return Inches(size)
        elif self.unit == Cm:
            return Cm(size)
        else:
            return Pt(size)
    def newer(self,filename:str):
        self.ppt = pptx.Presentation()
        self.filename = filename
        return self
    def loader(self, filename: str):
        #"./example_data/fake-power-point.pptx"
        self.ppt = pptx.Presentation(filename)
        self.filename = filename
        return self
    def lazy_load(self):
        if not self.ppt:
            raise Exception("没有正确指定ppt loader!")
        for slide in self.ppt.slides:
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text = text + text_frame.text + "\n"
            yield Document(text)
    def load(self):
        docs = []
        for doc in self.lazy_load():
            docs.append(doc)
        return docs

    def load_and_split(self, docx_file, chunk_size=1000,chunk_overlap=0):
        loader = self.loader(docx_file)
        docs = loader.load()
        spliter:RecursiveCharacterTextSplitter = self.langchainLib.get_textsplitter(chunk_size=chunk_size,chunk_overlap=chunk_overlap)
        splited_docs = spliter.split_documents(docs)
        return splited_docs
    def add_slide(self,slide_layout=None):
        if not slide_layout:
            slide_layout = self.ppt.slide_layouts[1]
        else:
            slide_layout = self.ppt.slide_layouts[slide_layout]
        slide = self.ppt.slides.add_slide(slide_layout)
        self.slides.append(slide)
        return self
    def set_title(self,title,subtitle=None,slide_idx=None):
        if slide_idx:
            slide = self.slides[slide_idx]
        else:
            slide = self.slides[-1]
        shapes = slide.shapes
        shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle
        return self
    def add_text(self,text,left,top,width,height,
                 word_wrap=False,
                 name=None,
                 bold=None,
                 italic=None,
                 size=None,
                 href=None,
                 alignment=None,
                 color=None,
                 slide_idx=None):
        if slide_idx:
            slide = self.slides[slide_idx]
        else:
            slide = self.slides[-1]
        left,top,width,height = self._get_sizes(left,top,width,height)
        shapes = slide.shapes
        txBox = shapes.add_textbox(left,top,width,height)
        tf = txBox.text_frame
        if word_wrap:
            tf.word_wrap = True
        if isinstance(text,list):
            tf.text = text[0]
            for rest_text in text[1:]:
                self.add_text_paragraph(tf,rest_text,
                                        name=name,
                                        bold=bold,
                                        italic=italic,
                                        size=size,
                                        href=href,
                                        alignment=alignment,
                                        color=color,
                                        level=1)
        else:
            tf.text = text
        return tf
    def add_text_paragraph(self,tf,text,
                           name=None,bold=False,italic=False,size=20,
                           href=None,
                           alignment=None,
                           color:Tuple[int,int,int]=None,level=0):
        p=tf.add_paragraph()
        p.text = text
        if name:
            p.font.name = name
        p.font.blod = bold
        p.font.italic = italic
        p.font.size = Pt(size)
        if alignment:
            if alignment.lower()=="left":
                p.alignment = PP_ALIGN.LEFT
            elif alignment.lower()=="right":
                p.alignment = PP_ALIGN.RIGHT
            elif alignment.lower()=="center":
                p.alignment = PP_ALIGN.CENTER
        if color:
            p.font.color.rgb = RGBColor(*color)
        if href:
            p.hyperlink.address = href
        p.level = level
        return p
    def add_image(self,image_path,left,top,width,height,slide_idx=None):
        if slide_idx:
            slide = self.slides[slide_idx]
        else:
            slide = self.slides[-1]
        left,top,width,height = self._get_sizes(left,top,width,height)
        shapes = slide.shapes
        image = shapes.add_picture(image_path,left,top,width,height)
        return image
    def add_chart(self,chart_type:Literal['pie','column','scatter','bubble'],
                  data:List[Dict[str,Dict]] | List[Dict[str,List[Tuple[int,int]]]] | List[Dict[str,List[Tuple[int,int,int]]]],
                  x,y,cx,cy,
                  smooth:List[bool]=None,
                  label = False,
                  label_size = None,
                  label_color: Tuple[int,int,int] = None,
                  legend = True,
                  legend_position:Literal['top','left','right','bottom'] = None,
                  slide_idx=None):
        if slide_idx:
            slide = self.slides[slide_idx]
        else:
            slide = self.slides[-1]
        x,y,cx,cy = self._get_sizes(x,y,cx,cy)
        shapes = slide.shapes
        if chart_type.lower() == 'pie':
            chart_type = XL_CHART_TYPE.PIE
        elif chart_type.lower() == 'column':
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        elif chart_type.lower() == 'scatter':
            chart_type = XL_CHART_TYPE.XY_SCATTER
        elif chart_type.lower() == 'bubble':
            chart_type = XL_CHART_TYPE.BUBBLE
        elif chart_type.lower() == 'line':
            chart_type = XL_CHART_TYPE.LINE
        else:
            chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

        
        if chart_type == XL_CHART_TYPE.XY_SCATTER:
            chart_data = XyChartData()
            for idx,item in enumerate(data):
                s = chart_data.add_series(list(item.keys())[0])
                for value in list(item.values())[0]:
                    s.add_data_point(*value)
        elif chart_type == XL_CHART_TYPE.BUBBLE:
            chart_data = BubbleChartData()
            for idx,item in enumerate(data):
                s = chart_data.add_series(list(item.keys())[0])
                for value in list(item.values())[0]:
                    s.add_data_point(*value)
        else:
            chart_data = ChartData()
            categories = list(list(data[0].values())[0].keys())
            chart_data.categories = categories
            for idx,item in enumerate(data):
                chart_data.add_series(list(item.keys())[0],list(list(item.values())[0].values()))
        chart_frame = shapes.add_chart(chart_type,x,y,cx,cy,chart_data)
        chart = chart_frame.chart
        if smooth and chart_type.lower == 'line' and len(smooth)==len(data):
            for idx,is_smooth in enumerate(smooth):
                chart_data.series[idx].smooth = is_smooth
        if label:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            if label_size:
                data_labels.font.size = self._get_size(label_size)
            if label_color:
                data_labels.font.color.rgb = RGBColor(*label_color)
        if legend:
            chart.has_legend = True
            if legend_position:
                if legend_position.lower() == 'left':
                    chart.legend.position = XL_LEGEND_POSITION.LEFT
                elif legend_position.lower() == 'right':
                    chart.legend.position = XL_LEGEND_POSITION.RIGHT
                elif legend_position.lower() == 'top':
                    chart.legend.position = XL_LEGEND_POSITION.TOP
                elif legend_position.lower() == 'bottom':
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        return chart
    
    def add_shape(self,shape_type,left,top,width,height,text=None,
                  fill_type:Literal["solid","transparent"]|None=None,
                  fill_color=None,
                  line_color=None,
                  line_brightness=None,
                  line_width=None,
                  rotation=None,
                  slide_idx=None) -> Shape:
        if slide_idx:
            slide = self.slides[slide_idx]
        else:
            slide = self.slides[-1]
        left,top,width,height = self._get_sizes(left,top,width,height)
        shapes = slide.shapes
        shape = shapes.add_shape(shape_type,left,top,width,height)
        shape.text = text
        if fill_type:
            if fill_type=="solid":
                shape.fill.solid()
            elif fill_type=="transparent":
                shape.fill.background()
        if fill_color and fill_type!="transparent":
            shape.fill.color.rgb = RGBColor(*fill_color)
        if line_color: 
            shape.line.color.rgb = RGBColor(*line_color)
        if line_brightness:
            shape.line.brightness = line_brightness
        if line_width:
            shape.line.width = self._get_size(line_width)
        if rotation:
            shape.rotation = rotation
        print(type(shape))
        return shape
    def get_shape_sizes(self,shape:Shape,unit=None):
        if not unit:
            unit = self.unit
        if unit==Inches:
            return shape.left.inches,shape.top.inches,shape.width.inches,shape.height.inches
        elif unit == Cm:
            return shape.left.cm,shape.top.cm,shape.width.cm,shape.height.cm
        else:
            return shape.left.pt,shape.top.pt,shape.width.pt,shape.height.pt 
    def add_table(self,contents:List[Dict],left,top,width,height,slide_idx=None,with_header=False):
        if slide_idx:
            slide = self.slides[slide_idx]
        else:
            slide = self.slides[-1]
        left,top,width,height = self._get_sizes(left,top,width,height)
        shapes = slide.shapes
        new_contents = contents.copy()
        if len(new_contents)==0:
            raise Exception("表格应至少包含一行一列!")
        if with_header:
            header = {}
            for key in new_contents[0]:
                header[key] = key
            new_contents.insert(0,header)
        
        rows = len(new_contents)
        cols = len(new_contents[0].keys())
        
        table = shapes.add_table(rows,cols,left,top,width,height).table
        for idx,row in enumerate(new_contents):
            for jdx,col_key in enumerate(row):
                table.cell(idx,jdx).text = str(new_contents[idx][col_key])
        return table
    def get_slide_master(self):
        slide_master = self.ppt.slide_master
        return slide_master
    def save(self):
        self.ppt.save(self.filename)
