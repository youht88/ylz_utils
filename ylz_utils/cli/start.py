from ylz_utils import Color, StringLib
from ylz_utils import FileLib
from ylz_utils import LangchainLib

from pydantic import BaseModel,Field
from typing import Literal,List
from pptx.util import Inches,Cm,Pt
from pptx.enum.shapes import MSO_SHAPE

from langchain.prompts import PromptTemplate,ChatPromptTemplate
from langchain.output_parsers import OutputFixingParser

from langchain_core.tools import tool
from langchain_core.runnables import RunnablePassthrough,RunnableLambda,RunnableParallel
from langchain_core.messages import HumanMessage

from ylz_utils.cli.chat import start_chat
from ylz_utils.cli.graph import start_graph
from ylz_utils.cli.rag import start_loader,start_rag

def __agent(langchainLib:LangchainLib,args):
    llm_key = args.llm_key
    llm_model = args.llm_model
    message = args.message if args.message else "厦门今天天气怎么样?"
    llm = langchainLib.get_llm(llm_key,llm_model)
    prompt = langchainLib.get_prompt()
    tools = []
    tavily_tool = langchainLib.toolLib.web_search.get_tool()
    tools.append(tavily_tool)
    # langchain_docs_vectorestore = langchainLib.vectorstoreLib.faiss.load("langchain_docs.db")
    # langchain_docs_retriever = langchain_docs_vectorestore.as_retriever()
    # rag_tool = langchainLib.toolLib.rag_search.get_tool(langchain_docs_retriever,"langchain docs","langchain的文档")
    # tools.append(rag_tool)
    
    agent =  langchainLib.get_agent(llm,tools)
    res = agent.stream({"messages":[("user",message)]})
    for chunk in res:
        print(chunk , end="")
    print("\n\n",langchainLib.get_llm(full=True),"\n")



def __llm_test(langchainLib:LangchainLib,args):
    llm_key = args.llm_key
    llm_model = args.llm_model
    message = args.message
    llm = langchainLib.get_llm(llm_key,llm_model)
    res = llm.invoke(message)
    print("res:",res)
    print("llms:",[(item["type"],item["api_key"],item["used"]) for item in langchainLib.get_llm(full=True)])

def __outputParser_test(langchainLib:LangchainLib,args):
    message = args.message
    llm_key = args.llm_key
    class Food(BaseModel):
        name: str = Field(description="name of the food")
        place: str = Field(defualt="未指定",description="where to eat food?",examples=["家","公司","体育馆"])
        calories: float|None = Field(title="卡路里(千卡)",description="食物热量")
        amount: tuple[float,str] = Field(description="tuple of the value and unit of the food amount")
        time: str = Field(default = "未指定",description="the time of eating food,确保每个食物都有对应的时间，如果没有指定则为空字符串")
        health: Literal["健康","不健康","未知"]|None = Field(description="根据常识判断食物是否为健康食品",examples=["健康","不健康","未知"])
    class Foods(BaseModel):
        foods: List[Food]

    # Set up a parser + inject instructions into the prompt template.
    parser = langchainLib.get_outputParser(Foods)

    prompt1 = PromptTemplate(
        template="用中文分析句子的内容。如果没有指定食物热量则根据食物的名称和数量进行估计。判断食物是否为健康\n{format_instructions}\n{command}\n",
        input_variables=["command"],
        partial_variables={"format_instructions": parser.get_format_instructions()},
    )


    # And a query intended to prompt a language model to populate the data structure.
    fixed_parser = OutputFixingParser.from_llm(parser=parser, llm=langchainLib.get_llm(llm_key))
    # retry_parser = RetryOutputParser.from_llm(parser=parser, llm=langchainLib.get_llm())

    prompt2 = langchainLib.get_prompt(
            system_prompt = "用中文分析句子的内容。如果没有指定食物热量则根据食物的名称和数量进行估计。同时判断食物是否为健康",
            human_keys={"command":"输入的语句"},
            outputParser = parser,
            use_chat = False
        )

    #chain = prompt2 | langchainLib.get_llm() | fixed_parser
    
    chain = prompt2 | langchainLib.get_llm(llm_key) | parser
    
    # retry_chain = RunnableParallel(
    #     completion=chain, prompt_value=prompt1,
    #     ) | RunnableLambda(lambda x:retry_parser.parse_with_prompt(**x))

    #output = chain.invoke({"command": "我今天在家吃了3个麦当劳炸鸡和一个焦糖布丁，昨天8点在电影院吃了一盘20千卡的西红柿炒蛋"})
    output = chain.invoke({"command": message})
    #retry_parser.parse_with_prompt(output['completion'], output['prompt_value'])
    #retry_parser.parse_with_prompt(output['completion'],output['prompt_value'])
    print("\noutput parser:",output)

    print("llms:",[(item["type"],item["api_key"],item["used"]) for item in langchainLib.get_llm(full=True)])


def __runnalble_test(langchainLib:LangchainLib):
    runnable1 = RunnableParallel(
    passed=RunnablePassthrough(),
    modified=lambda x: x["num"] + 1,
    other = RunnableParallel(
        passed=RunnablePassthrough(),
        modified=lambda x: x["num"] * 2,
        )
    )
    #runnable1.invoke({"num":1})
    def f(x):
        return x['modified']*100

    runnable2 = runnable1 | RunnableParallel(
        {"origin": RunnablePassthrough(),
        "new1": lambda x: x["other"]["passed"],
        "new2": RunnableLambda(f)}
        )

    #for trunck in runnable2.stream({"num": 1}):
    #  print(trunck)
    res = runnable2.invoke({"num": 1})
    print("\nrunnalbe:",res)

def __prompt_test(langchainLib:LangchainLib,args):
    llm_key = args.llm_key
    llm_model = args.llm_model
    prompt = langchainLib.get_prompt(f"你对日历时间非常精通",human_keys={"context":"关联的上下文","question":"问题是"},use_chat = False)
    prompt.partial(context="我的名字叫小美")   # 这个为什么不起作用?
    print("!!!!",langchainLib.llmLib.default_llm_key,llm_key)
    llm = langchainLib.get_llm(llm_key,llm_model)
    print(llm)
    outputParser = langchainLib.get_outputParser()
    chain = prompt | llm | outputParser

    res = chain.invoke({"question":"我叫什么名字，今天礼拜几","context":"我是海涛"})
    print(res)
    res = chain.stream({"question":"用我的名字写一首诗歌","context":"我的名字叫小美"})
    for chunk in res:
        print(chunk,end="")
    
    prompt = PromptTemplate.from_template("用中文回答：{topic} 的{what}是多少")
    llm = langchainLib.get_llm("LLM.GROQ")
    print(llm)
    chain = prompt | llm | outputParser
    promise = chain.batch([{"topic":"中国","what":"人口"},{"topic":"美国","what":"国土面积"},{"topic":"新加坡","what":"大学"}])
    print(promise)
    print("llms:",[(item["type"],item["api_key"],item["used"]) for item in langchainLib.get_llm(full=True)])


def __tools_test(langchainLib:LangchainLib,args):
    tool = langchainLib.toolLib.python_repl.get_tool()
    res = tool.invoke(
"""
print(3+2)
""")
    print("python repl=",res)
    tool = langchainLib.toolLib.wolfram_alpha.get_tool()
    res = tool.run("what is the square root of 25?")
    print(res)
    return 
    # tool: TavilySearchResults = langchainLib.get_search_tool("TAVILY")
    # prompt = langchainLib.get_prompt(is_chat=False,human_keys={"context":"关联的上下文是:","question":"问题是:"})
    # res = tool.invoke("易联众现在股价是多少？")
    # print(res)
    # llm = langchainLib.get_llm("LLM.DEEPSEEK")
    # chain = RunnableParallel({
    #     "question": RunnablePassthrough(),
    #     "context": tool
    # }) | prompt | llm 
    # res = chain.invoke("易联众现在股价是多少？")
    # print(res)

    # print("#"*50)
    class Output(BaseModel):
        name:str = Field(description= "姓名")
        age:int = Field(description= "年龄")

    langchainLib.add_plugins()
    llm_key = args.llm_key
    llm = langchainLib.get_llm(llm_key)
    print("llms:",[(item["type"],item["api_key"],item["used"]) for item in langchainLib.get_llm(full=True)])

    outputParser = langchainLib.get_outputParser(Output,fix=True,llm = llm,retry=3)
    prompt = langchainLib.get_prompt(human_keys = {"ask":"问题"},use_chat=False)
    search = langchainLib.get_search_tool("TAVILY")
    tools = [search]
    #chain  = prompt | (lambda x: x.messages[-1].content) | search | (lambda x:'```json\n{"name":中国,"age":"20"}') | outputParser
    #chain  = prompt | llm | outputParser
    chain = prompt | RunnableLambda(lambda x:x.text,name="抽取text") | search 
    #print("chain to_json = ",chain.to_json())
    #print("chain name=",chain.get_name())
    print("\nchain graph=",chain.get_graph(),"\n")
    res = chain.invoke({"ask": "北京人口多少"})
    print(type(res),res)
    #FileLib.writeFile("graph1.png",chain.get_graph().draw_mermaid_png(),mode="wb")

    # llm = langchainLib.get_llm("LLM.DEEPSEEK")
    # print("llm",llm.to_json())
    llm.bind_tools(tools)
    res = llm.invoke({"ask":"北京2024年人口多少"})
    print(res)

def start(args):
    langchainLib = LangchainLib()
    #langchainLib.add_plugins()
    if args.mode == "llm":
        StringLib.logging_in_box(f"\n{Color.YELLOW} 测试llm {Color.RESET}")
        __llm_test(langchainLib,args)
    elif args.mode == "prompt":
        StringLib.logging_in_box(f"\n{Color.YELLOW} 测试prompt {Color.RESET}")
        __prompt_test(langchainLib,args)    
    elif args.mode == 'loader':    
        StringLib.logging_in_box(f"\n{Color.YELLOW} 测试loader {Color.RESET}")
        start_loader(langchainLib,args)    
    elif args.mode == 'runnable':
        StringLib.logging_in_box(f"\n{Color.YELLOW} 测试runnable {Color.RESET}")
        __runnalble_test(langchainLib)
    elif args.mode == 'outputParser':    
        StringLib.logging_in_box(f"{Color.YELLOW} 测试outputParser {Color.RESET}")
        __outputParser_test(langchainLib,args)
    elif args.mode == 'rag':    
        StringLib.logging_in_box(f"\n{Color.YELLOW} 测试rag {Color.RESET}")
        start_rag(langchainLib,args)
    elif args.mode == 'tools':
        StringLib.logging_in_box(f"\n{Color.YELLOW} 测试tools {Color.RESET}")
        __tools_test(langchainLib,args)
    elif args.mode == 'graph':
        StringLib.logging_in_box(f"\n{Color.YELLOW} 测试graph {Color.RESET}")
        start_graph(langchainLib,args)
    elif args.mode == 'chat':
        start_chat(langchainLib,args)
    elif args.mode == 'agent':
        __agent(langchainLib,args)
    else:
        print(StringLib.yellow("args="),args)
        print(StringLib.yellow("llms--->:"),[(item["type"],item["api_key"],item["model"],item["used"]) for item in langchainLib.get_llm(full=True)])
        print(StringLib.yellow("embeddings---->:"),[(item["type"],item["api_key"],item["model"],item["used"]) for item in langchainLib.get_embedding(full=True)])
        #return
        agent = langchainLib.agentLib.get_full_agent(llm_key=args.llm_key,llm_model=args.llm_model)
        res = agent.invoke("今年的奥运会中国获得多少金牌？")
        return
        #langchainLib.ttsLib.tts_save("你好呀","tts.wav")
        #langchainLib.ttsLib.tts_play("祝你有愉快的一天")
        
        #loader = langchainLib.documentLib.pptx.loader("30335320.pptx")
        #docs = loader.load()
        #print(docs)
        
        # ## test elasticsearch vectorstore
        # langchainLib.vectorstoreLib.esLib.init_client(es_user="elastic")
        # langchainLib.vectorstoreLib.esLib.delete_store("langchain_index")
        # embed = langchainLib.embeddingLib.get_embedding('EMBEDDING.HF')
        # store = langchainLib.vectorstoreLib.esLib.get_store(embedding=embed)
        # ids = langchainLib.vectorstoreLib.esLib.create_from_texts(store,["我是一个科技工作者","树上有一只猴子"])
        # print(ids)
        # res = langchainLib.vectorstoreLib.esLib.search_with_score("我是谁？",store)
        # print(res)

        # embed = langchainLib.embeddingLib.get_embedding(fake_size=3)
        # print(embed.embed_documents(["hello world","I am fine!","我们正在学习langchain"]))
        
        # images = [
        #     {"image":"file://smoke.jpg"}
        # ]
        # res = langchainLib.get_llm("LLM.DASHSCOPE").invoke([HumanMessage(content=["图片中的文字是什么?"]+images)])
        
        image_message = {
                #"image": "file:///Users/youht/source/python/ylz_utils/tests/smoke.jpg",
                "image": "https://developer.qcloudimg.com/http-save/yehe-170434/6ac5c23a5efe1dfe00726f845eca9521.png",
            }
        text_message = {
            "text": "用中文详细解释这张图片的每一个细节,特别是每一个流程图的含义",
        }
        message = HumanMessage(content=[text_message, image_message])
        res = langchainLib.llmLib.get_llm(model="qwen-vl-max").invoke([message])
        print(StringLib.color("content=",style=["italic","bold","cyan","red"]),res.content)

        return 

        loader = langchainLib.documentLib.pptx.newer("test.pptx")
        loader.add_slide(0).set_title("Hello World","gogogo").add_text("youht",10,10,60,40)
        tab = [{"name":"youht","age":20},{"name":"jinli","age":10}] 
        tx = loader.add_slide(1).set_title("你好","step1").add_text("Step1",10,10,100,40)
        loader.add_text_paragraph(tx,"如何学习python",size=30,level=1)
        loader.add_text_paragraph(tx,"numpy",bold=True,size=20,level=2)
        loader.add_slide(2).set_title("你好","step1").add_table(tab,100,100,600,400,with_header = True)
        shape = loader.add_slide(8).set_title("Shape").add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,50,50,50,50,
                                                        line_brightness=0.5,
                                                        text="hello")
        print(loader.get_shape_sizes(shape))
        loader.add_slide().set_title("chart1").add_chart("pie",
                                                        [{"s1":{"a":1,"b":2,"c":3}},
                                                         {"s2":{"a":3,"b":2,"c":1}}],10,10,600,300)
        loader.add_slide().set_title("chart2").add_chart("bubble",
                                                        [{"s1":[(1,2,10),(4,6,3),(2,3,1)]},
                                                         {"s2":[(2,1,2),(8,2,3),(4,1,5)]}],10,10,600,300)

        loader.save()